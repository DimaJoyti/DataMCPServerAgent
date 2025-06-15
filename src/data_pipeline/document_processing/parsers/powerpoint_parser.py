"""
PowerPoint file parser for .pptx files.
"""

import logging
from pathlib import Path
from typing import Dict, List

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    HAS_PYTHON_PPTX = True
except ImportError:
    HAS_PYTHON_PPTX = False

from ..metadata.models import DocumentMetadata, DocumentType
from .base_parser import BaseParser, ParsedDocument


class PowerPointParser(BaseParser):
    """Parser for PowerPoint files (.pptx)."""

    def __init__(self):
        """Initialize PowerPoint parser."""
        super().__init__()
        self.logger = logging.getLogger(self.__class__.__name__)

        if not HAS_PYTHON_PPTX:
            raise ImportError(
                "python-pptx is required for PowerPoint parsing. Install with: pip install python-pptx"
            )

    def can_parse(self, file_path: Path) -> bool:
        """Check if this parser can handle the file."""
        return file_path.suffix.lower() == ".pptx"

    def parse(self, file_path: Path, **kwargs) -> ParsedDocument:
        """
        Parse PowerPoint file.

        Args:
            file_path: Path to PowerPoint file
            **kwargs: Additional parsing options
                - include_notes: Whether to include speaker notes (default: True)
                - include_slide_numbers: Whether to include slide numbers (default: True)
                - extract_tables: Whether to extract table content (default: True)
                - slide_separator: Separator between slides (default: '\n\n---\n\n')

        Returns:
            ParsedDocument: Parsed document
        """
        try:
            # Parse options
            include_notes = kwargs.get("include_notes", True)
            include_slide_numbers = kwargs.get("include_slide_numbers", True)
            extract_tables = kwargs.get("extract_tables", True)
            slide_separator = kwargs.get("slide_separator", "\n\n---\n\n")

            # Load presentation
            presentation = Presentation(file_path)

            # Extract content
            slides_content = []
            slides_metadata = []
            total_word_count = 0

            for slide_idx, slide in enumerate(presentation.slides, 1):
                slide_content, slide_meta = self._extract_slide_content(
                    slide,
                    slide_idx,
                    include_notes=include_notes,
                    include_slide_numbers=include_slide_numbers,
                    extract_tables=extract_tables,
                )

                slides_content.append(slide_content)
                slides_metadata.append(slide_meta)
                total_word_count += slide_meta.get("word_count", 0)

            # Combine all slides
            full_text = slide_separator.join(slides_content)

            # Extract presentation metadata
            pres_metadata = self._extract_presentation_metadata(presentation)

            # Create document metadata
            doc_metadata = DocumentMetadata(
                title=pres_metadata.get("title") or file_path.stem,
                document_type=DocumentType.PRESENTATION,
                file_path=str(file_path),
                file_size=file_path.stat().st_size,
                language="unknown",  # PowerPoint doesn't have inherent language detection
                page_count=len(presentation.slides),
                word_count=total_word_count,
                character_count=len(full_text),
                custom_metadata={
                    "total_slides": len(presentation.slides),
                    "slides_metadata": slides_metadata,
                    "presentation_metadata": pres_metadata,
                    "parser": "PowerPointParser",
                },
            )

            return ParsedDocument(
                text=full_text,
                metadata=doc_metadata,
                raw_content=presentation,  # Store original presentation object
                processing_time=0.0,  # Will be set by caller
            )

        except Exception as e:
            self.logger.error(f"Failed to parse PowerPoint file {file_path}: {e}")
            raise

    def _extract_slide_content(
        self,
        slide,
        slide_number: int,
        include_notes: bool = True,
        include_slide_numbers: bool = True,
        extract_tables: bool = True,
    ) -> tuple[str, Dict]:
        """Extract content from a single slide."""
        content_parts = []
        metadata = {
            "slide_number": slide_number,
            "shapes_count": len(slide.shapes),
            "has_title": False,
            "has_content": False,
            "has_notes": False,
            "has_tables": False,
            "has_images": False,
            "word_count": 0,
        }

        # Add slide number if requested
        if include_slide_numbers:
            content_parts.append(f"Slide {slide_number}")

        # Extract text from shapes
        slide_text_parts = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                slide_text_parts.append(text)

                # Check if this is likely a title
                if shape == slide.shapes[0] or len(text) < 100:
                    metadata["has_title"] = True
                else:
                    metadata["has_content"] = True

            # Handle tables
            if extract_tables and shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_text = self._extract_table_content(shape.table)
                if table_text:
                    slide_text_parts.append(table_text)
                    metadata["has_tables"] = True

            # Check for images
            if shape.shape_type in [MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.MEDIA]:
                metadata["has_images"] = True

        # Add slide content
        if slide_text_parts:
            content_parts.extend(slide_text_parts)

        # Extract speaker notes
        if include_notes and slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                content_parts.append(f"Speaker Notes: {notes_text}")
                metadata["has_notes"] = True

        # Combine slide content
        slide_content = "\n".join(content_parts)
        metadata["word_count"] = len(slide_content.split())

        return slide_content, metadata

    def _extract_table_content(self, table) -> str:
        """Extract content from a table."""
        try:
            table_rows = []

            for row in table.rows:
                row_cells = []
                for cell in row.cells:
                    cell_text = cell.text.strip() if cell.text else ""
                    row_cells.append(cell_text)

                if any(row_cells):  # Only add non-empty rows
                    table_rows.append(" | ".join(row_cells))

            if table_rows:
                return "Table:\n" + "\n".join(table_rows)

        except Exception as e:
            self.logger.warning(f"Failed to extract table content: {e}")

        return ""

    def _extract_presentation_metadata(self, presentation) -> Dict:
        """Extract metadata from presentation."""
        metadata = {}

        try:
            # Core properties
            core_props = presentation.core_properties
            if core_props:
                metadata.update(
                    {
                        "title": core_props.title,
                        "author": core_props.author,
                        "subject": core_props.subject,
                        "keywords": core_props.keywords,
                        "comments": core_props.comments,
                        "category": core_props.category,
                        "created": core_props.created.isoformat() if core_props.created else None,
                        "modified": (
                            core_props.modified.isoformat() if core_props.modified else None
                        ),
                        "last_modified_by": core_props.last_modified_by,
                        "revision": core_props.revision,
                        "version": core_props.version,
                        "language": core_props.language,
                        "content_status": core_props.content_status,
                    }
                )

            # Slide dimensions
            slide_width = presentation.slide_width
            slide_height = presentation.slide_height
            metadata.update(
                {
                    "slide_width": slide_width,
                    "slide_height": slide_height,
                    "slide_aspect_ratio": (
                        round(slide_width / slide_height, 2) if slide_height > 0 else None
                    ),
                }
            )

            # Slide layouts info
            layouts_info = []
            for layout in presentation.slide_layouts:
                layout_info = {"name": layout.name, "placeholders_count": len(layout.placeholders)}
                layouts_info.append(layout_info)

            metadata["slide_layouts"] = layouts_info

            # Master slides info
            masters_info = []
            for master in presentation.slide_masters:
                master_info = {"layouts_count": len(master.slide_layouts)}
                masters_info.append(master_info)

            metadata["slide_masters"] = masters_info

        except Exception as e:
            self.logger.warning(f"Failed to extract presentation metadata: {e}")

        return metadata

    def extract_metadata(self, file_path: Path) -> Dict:
        """Extract metadata from PowerPoint file."""
        try:
            presentation = Presentation(file_path)
            return self._extract_presentation_metadata(presentation)
        except Exception as e:
            self.logger.warning(f"Failed to extract PowerPoint metadata: {e}")
            return {}

    def get_supported_extensions(self) -> List[str]:
        """Get list of supported file extensions."""
        return [".pptx"]

    def get_parser_info(self) -> Dict:
        """Get information about this parser."""
        return {
            "name": "PowerPointParser",
            "description": "Parser for Microsoft PowerPoint files",
            "supported_extensions": self.get_supported_extensions(),
            "features": [
                "Slide content extraction",
                "Speaker notes extraction",
                "Table content extraction",
                "Image detection",
                "Slide metadata",
                "Presentation properties",
            ],
            "dependencies": {"python-pptx": HAS_PYTHON_PPTX},
            "limitations": [
                "Only supports .pptx format (not .ppt)",
                "Image content is not extracted (only detected)",
                "Complex shapes may not be fully parsed",
            ],
        }
