"""
Export tools for the Research Assistant.

This module provides tools for exporting research results in various formats,
including PDF, DOCX, Markdown, HTML, and presentation formats.
"""

import json
from datetime import datetime
from typing import Dict

from langchain.tools import Tool


class MarkdownExporter:
    """Tool for exporting research results to Markdown."""

    def export_to_markdown(self, research_data: Dict) -> str:
        """
        Export research data to Markdown format.

        Args:
            research_data: Dictionary containing research data

        Returns:
            Markdown string
        """
        # Extract research data
        topic = research_data.get("topic", "Research Topic")
        summary = research_data.get("summary", "")
        sources = research_data.get("sources", [])
        tools_used = research_data.get("tools_used", [])

        # Format the Markdown
        markdown = f"# {topic}\n\n"
        markdown += f"## Summary\n\n{summary}\n\n"

        # Add sources
        markdown += "## Sources\n\n"
        for i, source in enumerate(sources, 1):
            if isinstance(source, dict):
                title = source.get("title", f"Source {i}")
                url = source.get("url", "")
                authors = source.get("authors", [])

                markdown += f"{i}. **{title}**\n"
                if authors:
                    markdown += f"   Authors: {', '.join(authors)}\n"
                if url:
                    markdown += f"   URL: {url}\n"
                markdown += "\n"
            else:
                markdown += f"{i}. {source}\n\n"

        # Add tools used
        markdown += "## Tools Used\n\n"
        for tool in tools_used:
            markdown += f"- {tool}\n"

        # Add timestamp
        markdown += (
            f"\n\n---\nGenerated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
        )

        return markdown

    def save_markdown_to_file(self, markdown: str, filename: str) -> str:
        """
        Save Markdown to a file.

        Args:
            markdown: Markdown string
            filename: Output filename

        Returns:
            Success message
        """
        # Ensure the filename has the .md extension
        if not filename.endswith(".md"):
            filename += ".md"

        # Write the Markdown to the file
        with open(filename, "w", encoding="utf-8") as f:
            f.write(markdown)

        return f"Markdown exported to {filename}"

    def run(self, input_str: str) -> str:
        """
        Run the Markdown exporter.

        Args:
            input_str: JSON string containing research data and filename

        Returns:
            Success message
        """
        try:
            data = json.loads(input_str)
            research_data = data.get("research_data", {})
            filename = data.get(
                "filename", f"research_{datetime.now().strftime('%Y%m%d%H%M%S')}.md"
            )

            markdown = self.export_to_markdown(research_data)
            return self.save_markdown_to_file(markdown, filename)
        except Exception as e:
            return f"Error exporting to Markdown: {str(e)}"


class HTMLExporter:
    """Tool for exporting research results to HTML."""

    def export_to_html(self, research_data: Dict) -> str:
        """
        Export research data to HTML format.

        Args:
            research_data: Dictionary containing research data

        Returns:
            HTML string
        """
        # Extract research data
        topic = research_data.get("topic", "Research Topic")
        summary = research_data.get("summary", "")
        sources = research_data.get("sources", [])
        tools_used = research_data.get("tools_used", [])

        # Format the HTML
        html = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{topic}</title>
    <style>
        body {{ font-family: Arial, sans-serif; line-height: 1.6; margin: 0; padding: 20px; max-width: 800px; margin: 0 auto; }}
        h1 {{ color: #333; }}
        h2 {{ color: #444; margin-top: 20px; }}
        .source {{ margin-bottom: 15px; }}
        .tools {{ margin-top: 20px; }}
        .footer {{ margin-top: 30px; font-size: 0.8em; color: #777; border-top: 1px solid #eee; padding-top: 10px; }}
    </style>
</head>
<body>
    <h1>{topic}</h1>

    <h2>Summary</h2>
    <p>{summary}</p>

    <h2>Sources</h2>
    <div class="sources">
"""

        # Add sources
        for i, source in enumerate(sources, 1):
            html += '        <div class="source">\n'
            if isinstance(source, dict):
                title = source.get("title", f"Source {i}")
                url = source.get("url", "")
                authors = source.get("authors", [])

                html += f"            <p><strong>{i}. {title}</strong></p>\n"
                if authors:
                    html += f"            <p>Authors: {', '.join(authors)}</p>\n"
                if url:
                    html += f'            <p>URL: <a href="{url}" target="_blank">{url}</a></p>\n'
            else:
                html += f"            <p>{i}. {source}</p>\n"
            html += "        </div>\n"

        # Add tools used
        html += """    </div>

    <h2>Tools Used</h2>
    <ul class="tools">
"""

        for tool in tools_used:
            html += f"        <li>{tool}</li>\n"

        # Add timestamp
        html += f"""    </ul>

    <div class="footer">
        <p>Generated on {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}</p>
    </div>
</body>
</html>
"""

        return html

    def save_html_to_file(self, html: str, filename: str) -> str:
        """
        Save HTML to a file.

        Args:
            html: HTML string
            filename: Output filename

        Returns:
            Success message
        """
        # Ensure the filename has the .html extension
        if not filename.endswith(".html"):
            filename += ".html"

        # Write the HTML to the file
        with open(filename, "w", encoding="utf-8") as f:
            f.write(html)

        return f"HTML exported to {filename}"

    def run(self, input_str: str) -> str:
        """
        Run the HTML exporter.

        Args:
            input_str: JSON string containing research data and filename

        Returns:
            Success message
        """
        try:
            data = json.loads(input_str)
            research_data = data.get("research_data", {})
            filename = data.get(
                "filename", f"research_{datetime.now().strftime('%Y%m%d%H%M%S')}.html"
            )

            html = self.export_to_html(research_data)
            return self.save_html_to_file(html, filename)
        except Exception as e:
            return f"Error exporting to HTML: {str(e)}"


class PDFExporter:
    """Tool for exporting research results to PDF."""

    def export_to_pdf(self, research_data: Dict, filename: str) -> str:
        """
        Export research data to PDF format.

        Args:
            research_data: Dictionary containing research data
            filename: Output filename

        Returns:
            Success message
        """
        try:
            # Import the FPDF library
            try:
                from fpdf import FPDF

                FPDF_AVAILABLE = True
            except ImportError:
                FPDF_AVAILABLE = False
                print("Warning: FPDF library not available. Using mock PDF generation.")

            # Ensure the filename has the .pdf extension
            if not filename.endswith(".pdf"):
                filename += ".pdf"

            if FPDF_AVAILABLE:
                # Extract research data
                topic = research_data.get("topic", "Research Topic")
                summary = research_data.get("summary", "")
                sources = research_data.get("sources", [])
                tools_used = research_data.get("tools_used", [])

                # Create PDF object
                pdf = FPDF()
                pdf.add_page()

                # Set font
                pdf.set_font("Arial", "B", 16)

                # Add title
                pdf.cell(0, 10, topic, 0, 1, "C")
                pdf.ln(10)

                # Add summary
                pdf.set_font("Arial", "B", 12)
                pdf.cell(0, 10, "Summary", 0, 1)
                pdf.set_font("Arial", "", 12)

                # Split summary into lines to fit the page width
                pdf.multi_cell(0, 10, summary)
                pdf.ln(10)

                # Add sources
                pdf.set_font("Arial", "B", 12)
                pdf.cell(0, 10, "Sources", 0, 1)
                pdf.set_font("Arial", "", 12)

                for i, source in enumerate(sources, 1):
                    if isinstance(source, dict):
                        title = source.get("title", f"Source {i}")
                        url = source.get("url", "")
                        authors = source.get("authors", [])

                        pdf.set_font("Arial", "B", 12)
                        pdf.cell(0, 10, f"{i}. {title}", 0, 1)
                        pdf.set_font("Arial", "", 12)

                        if authors:
                            pdf.cell(0, 10, f"Authors: {', '.join(authors)}", 0, 1)
                        if url:
                            pdf.cell(0, 10, f"URL: {url}", 0, 1)
                    else:
                        pdf.cell(0, 10, f"{i}. {source}", 0, 1)

                    pdf.ln(5)

                # Add tools used
                pdf.set_font("Arial", "B", 12)
                pdf.cell(0, 10, "Tools Used", 0, 1)
                pdf.set_font("Arial", "", 12)

                for tool in tools_used:
                    pdf.cell(0, 10, f"- {tool}", 0, 1)

                # Add timestamp
                pdf.ln(10)
                pdf.set_font("Arial", "I", 10)
                pdf.cell(
                    0,
                    10,
                    f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
                    0,
                    1,
                )

                # Save the PDF
                pdf.output(filename)

                return f"PDF exported to {filename}"
            else:
                # Fall back to mock implementation
                # First, generate HTML
                html_exporter = HTMLExporter()
                html = html_exporter.export_to_html(research_data)

                # Mock PDF generation
                with open(filename, "w", encoding="utf-8") as f:
                    f.write(f"PDF MOCK CONTENT\n\n{html}")

                return f"PDF exported to {filename} (mock implementation)"
        except Exception as e:
            return f"Error exporting to PDF: {str(e)}"

    def run(self, input_str: str) -> str:
        """
        Run the PDF exporter.

        Args:
            input_str: JSON string containing research data and filename

        Returns:
            Success message
        """
        try:
            data = json.loads(input_str)
            research_data = data.get("research_data", {})
            filename = data.get(
                "filename", f"research_{datetime.now().strftime('%Y%m%d%H%M%S')}.pdf"
            )

            return self.export_to_pdf(research_data, filename)
        except Exception as e:
            return f"Error exporting to PDF: {str(e)}"


class DOCXExporter:
    """Tool for exporting research results to DOCX."""

    def export_to_docx(self, research_data: Dict, filename: str) -> str:
        """
        Export research data to DOCX format.

        Args:
            research_data: Dictionary containing research data
            filename: Output filename

        Returns:
            Success message
        """
        try:
            # Import the python-docx library
            try:
                import docx
                from docx.shared import Inches, Pt

                DOCX_AVAILABLE = True
            except ImportError:
                DOCX_AVAILABLE = False
                print(
                    "Warning: python-docx library not available. Using mock DOCX generation."
                )

            # Ensure the filename has the .docx extension
            if not filename.endswith(".docx"):
                filename += ".docx"

            if DOCX_AVAILABLE:
                # Extract research data
                topic = research_data.get("topic", "Research Topic")
                summary = research_data.get("summary", "")
                sources = research_data.get("sources", [])
                tools_used = research_data.get("tools_used", [])

                # Create a new document
                doc = docx.Document()

                # Add title
                doc.add_heading(topic, 0)

                # Add summary
                doc.add_heading("Summary", level=1)
                doc.add_paragraph(summary)

                # Add sources
                doc.add_heading("Sources", level=1)

                for i, source in enumerate(sources, 1):
                    if isinstance(source, dict):
                        title = source.get("title", f"Source {i}")
                        url = source.get("url", "")
                        authors = source.get("authors", [])

                        p = doc.add_paragraph()
                        p.add_run(f"{i}. {title}").bold = True

                        if authors:
                            doc.add_paragraph(f"Authors: {', '.join(authors)}")
                        if url:
                            doc.add_paragraph(f"URL: {url}")
                    else:
                        doc.add_paragraph(f"{i}. {source}")

                    doc.add_paragraph()

                # Add tools used
                doc.add_heading("Tools Used", level=1)

                for tool in tools_used:
                    doc.add_paragraph(f"- {tool}", style="List Bullet")

                # Add timestamp
                doc.add_paragraph()
                p = doc.add_paragraph(
                    f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
                )
                p.style = "Subtitle"

                # Save the document
                doc.save(filename)

                return f"DOCX exported to {filename}"
            else:
                # Mock DOCX generation
                with open(filename, "w", encoding="utf-8") as f:
                    f.write(
                        f"DOCX MOCK CONTENT\n\nTopic: {research_data.get('topic', 'Research Topic')}"
                    )

                return f"DOCX exported to {filename} (mock implementation)"
        except Exception as e:
            return f"Error exporting to DOCX: {str(e)}"

    def run(self, input_str: str) -> str:
        """
        Run the DOCX exporter.

        Args:
            input_str: JSON string containing research data and filename

        Returns:
            Success message
        """
        try:
            data = json.loads(input_str)
            research_data = data.get("research_data", {})
            filename = data.get(
                "filename", f"research_{datetime.now().strftime('%Y%m%d%H%M%S')}.docx"
            )

            return self.export_to_docx(research_data, filename)
        except Exception as e:
            return f"Error exporting to DOCX: {str(e)}"


class PresentationExporter:
    """Tool for exporting research results to a presentation format."""

    def export_to_presentation(self, research_data: Dict, filename: str) -> str:
        """
        Export research data to a presentation format.

        Args:
            research_data: Dictionary containing research data
            filename: Output filename

        Returns:
            Success message
        """
        try:
            # Import the python-pptx library
            try:
                from pptx import Presentation
                from pptx.util import Inches, Pt

                PPTX_AVAILABLE = True
            except ImportError:
                PPTX_AVAILABLE = False
                print(
                    "Warning: python-pptx library not available. Using mock presentation generation."
                )

            # Ensure the filename has the .pptx extension
            if not filename.endswith(".pptx"):
                filename += ".pptx"

            if PPTX_AVAILABLE:
                # Extract research data
                topic = research_data.get("topic", "Research Topic")
                summary = research_data.get("summary", "")
                sources = research_data.get("sources", [])
                tools_used = research_data.get("tools_used", [])

                # Create a new presentation
                prs = Presentation()

                # Add title slide
                title_slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(title_slide_layout)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]

                title.text = topic
                subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d')}"

                # Add summary slide
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)
                title = slide.shapes.title
                content = slide.placeholders[1]

                title.text = "Summary"

                # Split summary into paragraphs
                paragraphs = summary.split("\n")
                text_frame = content.text_frame

                for i, paragraph in enumerate(paragraphs):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = paragraph

                # Add sources slide
                slide = prs.slides.add_slide(bullet_slide_layout)
                title = slide.shapes.title
                content = slide.placeholders[1]

                title.text = "Sources"
                text_frame = content.text_frame

                for i, source in enumerate(
                    sources[:5], 1
                ):  # Limit to 5 sources to fit on slide
                    if i == 1:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()

                    if isinstance(source, dict):
                        title_text = source.get("title", f"Source {i}")
                        authors = source.get("authors", [])

                        p.text = f"{title_text}"
                        if authors:
                            p.level = 0
                            p = text_frame.add_paragraph()
                            p.text = f"Authors: {', '.join(authors)}"
                            p.level = 1
                    else:
                        p.text = f"{source}"

                # Add tools used slide
                slide = prs.slides.add_slide(bullet_slide_layout)
                title = slide.shapes.title
                content = slide.placeholders[1]

                title.text = "Tools Used"
                text_frame = content.text_frame

                for i, tool in enumerate(tools_used):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = tool

                # Save the presentation
                prs.save(filename)

                return f"Presentation exported to {filename}"
            else:
                # Mock presentation generation
                with open(filename, "w", encoding="utf-8") as f:
                    f.write(
                        f"PPTX MOCK CONTENT\n\nTopic: {research_data.get('topic', 'Research Topic')}"
                    )

                return f"Presentation exported to {filename} (mock implementation)"
        except Exception as e:
            return f"Error exporting to presentation: {str(e)}"

    def run(self, input_str: str) -> str:
        """
        Run the presentation exporter.

        Args:
            input_str: JSON string containing research data and filename

        Returns:
            Success message
        """
        try:
            data = json.loads(input_str)
            research_data = data.get("research_data", {})
            filename = data.get(
                "filename", f"research_{datetime.now().strftime('%Y%m%d%H%M%S')}.pptx"
            )

            return self.export_to_presentation(research_data, filename)
        except Exception as e:
            return f"Error exporting to presentation: {str(e)}"


# Create tool instances
markdown_exporter = MarkdownExporter()
html_exporter = HTMLExporter()
pdf_exporter = PDFExporter()
docx_exporter = DOCXExporter()
presentation_exporter = PresentationExporter()

# Create LangChain tools
export_to_markdown_tool = Tool(
    name="export_to_markdown",
    func=markdown_exporter.run,
    description="Export research results to Markdown format. Input should be a JSON string with 'research_data' and 'filename' fields.",
)

export_to_html_tool = Tool(
    name="export_to_html",
    func=html_exporter.run,
    description="Export research results to HTML format. Input should be a JSON string with 'research_data' and 'filename' fields.",
)

export_to_pdf_tool = Tool(
    name="export_to_pdf",
    func=pdf_exporter.run,
    description="Export research results to PDF format. Input should be a JSON string with 'research_data' and 'filename' fields.",
)

export_to_docx_tool = Tool(
    name="export_to_docx",
    func=docx_exporter.run,
    description="Export research results to DOCX format. Input should be a JSON string with 'research_data' and 'filename' fields.",
)

export_to_presentation_tool = Tool(
    name="export_to_presentation",
    func=presentation_exporter.run,
    description="Export research results to a presentation format. Input should be a JSON string with 'research_data' and 'filename' fields.",
)
